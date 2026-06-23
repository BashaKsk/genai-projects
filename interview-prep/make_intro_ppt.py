"""Generate a one-slide PPT: 'Tell Me About Yourself' for Basha."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1A, 0x23, 0x7E)
ORANGE = RGBColor(0xFF, 0x99, 0x00)   # Amazon-ish accent
DARK = RGBColor(0x22, 0x22, 0x2A)
GRAY = RGBColor(0x5A, 0x5A, 0x6E)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# --- Top accent bar ---
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.25))
bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()

# --- Title ---
tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12), Inches(0.95))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Tell Me About Yourself"
r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "Basha  ·  Senior Full Stack Developer"
r2.font.size = Pt(16); r2.font.color.rgb = ORANGE

# --- Beats (CCBP pattern) ---
beats = [
    ("Background", "Bachelor's in Mechanical Engineering (2022) — fell in love with coding. "
                   "Today, a Senior Full Stack Developer."),
    ("Skills", "Python, JavaScript, React, Node.js — and GenAI with LangChain."),
    ("Goal", "Grow into building AI-powered, problem-solving products."),
    ("Why Amazon", "Drawn to customer-centric, AI-powered products — making people's lives easier."),
    ("Keystone Project", "Built a seafood traceability system: a Neo4j knowledge graph + a GenAI "
                         "chatbot. Never used a graph DB — learned it from scratch and shipped to production."),
    ("Closer", "I'd love to contribute my AI + full-stack skills here.   |   Outside work: cricket & gym."),
]

body = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(12), Inches(5.5))
bf = body.text_frame; bf.word_wrap = True
for i, (label, text) in enumerate(beats):
    para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
    para.space_after = Pt(10)
    lr = para.add_run(); lr.text = f"{label}:  "
    lr.font.bold = True; lr.font.size = Pt(17); lr.font.color.rgb = NAVY
    tr = para.add_run(); tr.text = text
    tr.font.size = Pt(15); tr.font.color.rgb = DARK

# --- Footer tip ---
foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12), Inches(0.4))
fp = foot.text_frame.paragraphs[0]
fr = fp.add_run()
fr.text = "Delivery: ~60–90 sec · slow down · lead with \"Senior Full Stack Developer\" · land the result."
fr.font.size = Pt(11); fr.font.italic = True; fr.font.color.rgb = GRAY

out = "/home/basha/Desktop/genai/interview-prep/tell-me-about-yourself.pptx"
prs.save(out)
print("Saved:", out)
